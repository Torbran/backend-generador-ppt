import sys # Asegúrate de que sys está importado si no lo estaba ya
print("DEBUG: === Inicio del script main.py ===", file=sys.stderr)
from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import List, Optional
import requests, tempfile, uuid, os
from pptx import Presentation
import sys # Importamos sys para escribir en la salida de errores (logs de Render)
from pptx.util import Inches # Importamos Inches para posicionar imágenes
import traceback # Importamos traceback para imprimir errores completos

# Obtener la ruta del directorio donde se encuentra este script (main.py)
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_FILENAME = "PLANTILLA_MARCADORES_AUTOMATIZADA.pptx"
# Construir la ruta completa al archivo de plantilla usando la ruta base
TEMPLATE_PATH = os.path.join(BASE_DIR, TEMPLATE_FILENAME)

# ---------- Models ---------- #
# Define la estructura de los datos que esperas recibir
class Camara(BaseModel):
    numero: int
    tipo: str
    ubicacion: str
    observaciones: str
    imagen_url: Optional[str] = None # La URL de la imagen es opcional

class Proyecto(BaseModel):
    nombre: str
    ubicacion: str
    fecha: str
    csv: str # Nombre del archivo CSV asociado (si aplica)
    mapa_url: str # URL de la imagen del mapa
    descripcion: str # Descripción del proyecto

class NVR(BaseModel):
    tipo: str
    direccion: str
    observaciones: Optional[str] = ""

class Carteles(BaseModel):
    barrio_protegido: int
    domiciliarios: int

class Cierre(BaseModel):
    circulo: str
    fecha_aprobacion: str

class InformeData(BaseModel):
    proyecto: Proyecto
    camaras: List[Camara] # Una lista de objetos Camara
    nvr: NVR
    carteles: Carteles
    cierre: Cierre

# ---------- Utils ---------- #
# Función para reemplazar texto en las formas de la presentación
def replace_text(shape, placeholder, value):
    if shape.has_text_frame:
        if placeholder in shape.text_frame.text:
            # Asegúrate de que 'value' es un string para evitar errores
            shape.text_frame.text = shape.text_frame.text.replace(placeholder, str(value))

# Función para descargar una imagen desde una URL a un archivo temporal
def download_image(url):
    try:
        print(f"DEBUG: Intentando descargar imagen desde URL: {url}", file=sys.stderr) # Log de inicio de descarga
        r = requests.get(url, timeout=15) # Aumentar timeout por si la descarga es lenta
        r.raise_for_status() # Lanza una excepción para códigos de estado HTTP de error (4xx o 5xx)

        # Usar tempfile para crear un archivo temporal seguro
        # Lo creamos en el directorio /tmp que es escribible en entornos como Render
        # Añadimos el sufijo correcto por si la librería de pptx lo necesita
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(url)[-1], dir='/tmp')
        tmp.write(r.content)
        tmp.close() # Cerrar el archivo para asegurar que se ha escrito

        print(f"DEBUG: Imagen descargada exitosamente a temporal: {tmp.name}", file=sys.stderr) # Log de éxito
        return tmp.name
    except requests.exceptions.RequestException as req_e:
        print(f"ERROR: Error de red o HTTP al descargar imagen desde {url}: {req_e}", file=sys.stderr) # Log de error de red/http
        traceback.print_exc(file=sys.stderr)
        return None
    except Exception as e:
        # === LOGGING DETALLADO DE CUALQUIER OTRO ERROR EN DESCARGA ===
        print(f"ERROR: Error inesperado al descargar imagen desde {url}: {e}", file=sys.stderr) # Log de otros errores
        traceback.print_exc(file=sys.stderr) # Imprimir rastreo completo
        return None # Devolver None si falla

# ---------- Main logic ---------- #
# Función principal para generar la presentación
def generar_presentacion(data: InformeData) -> str:
    # === VERIFICAR Y LOGUEAR LA RUTA DE LA PLANTILLA AL INICIO ===
    print(f"DEBUG: Verificando si el archivo de plantilla existe en la ruta: {TEMPLATE_PATH}", file=sys.stderr)
    if not os.path.exists(TEMPLATE_PATH):
        print(f"ERROR: Archivo de plantilla NO encontrado en la ruta esperada: {TEMPLATE_PATH}", file=sys.stderr)
        # Si la plantilla no existe, lanzamos una excepción clara que FastAPI capturará
        raise FileNotFoundError(f"El archivo de plantilla '{TEMPLATE_FILENAME}' no se encontró en la ruta del servidor: {TEMPLATE_PATH}")

    try:
        # === LOGUEAR ANTES DE CARGAR LA PRESENTACIÓN ===
        print(f"DEBUG: Archivo de plantilla encontrado. Procediendo a cargar la presentación desde {TEMPLATE_PATH}", file=sys.stderr)
        prs = Presentation(TEMPLATE_PATH) # Cargar la presentación desde la plantilla
        print(f"DEBUG: Plantilla cargada exitosamente.", file=sys.stderr) # Log de éxito al cargar


        # --- Procesar Slides Iniciales (0 y 1 según tu código original) ---
        # Rellenar campos dinámicos en las primeras diapositivas
        print("DEBUG: Rellenando campos en slides iniciales (0 y 1)...", file=sys.stderr)
        for slide in prs.slides[0:2]:
            for shape in slide.shapes:
                replace_text(shape, "{{nombre_proyecto}}", data.proyecto.nombre)
                replace_text(shape, "{{ubicacion}}", data.proyecto.ubicacion)
                replace_text(shape, "{{fecha}}", data.proyecto.fecha)
                replace_text(shape, "{{nombre_csv}}", data.proyecto.csv)
                # Asegúrate que {{calle_principal}} en la plantilla se mapea a data.proyecto.ubicacion o si tienes otro campo
                replace_text(shape, "{{calle_principal}}", data.proyecto.ubicacion) # O ajusta si tienes otro campo en Proyecto
                replace_text(shape, "{{descripcion}}", data.proyecto.descripcion)
        print("DEBUG: Campos de slides iniciales rellenados.", file=sys.stderr)

        # --- Procesar Slides de Cantidades/NVR ---
        # Asumo que los placeholders de cantidades y NVR están distribuidos en varios slides.
        # Este loop rellena esos placeholders donde sea que los encuentre en *todos* los slides.
        print("DEBUG: Rellenando campos de cantidades/NVR en todos los slides...", file=sys.stderr)
        for slide in prs.slides:
             for shape in slide.shapes:
                # Revisa que estos placeholders coincidan exactamente con los de tu plantilla
                replace_text(shape, "{{carteles_bp}}", str(data.carteles.barrio_protegido))
                replace_text(shape, "{{carteles_dom}}", str(data.carteles.domiciliarios))
                replace_text(shape, "{{nvr_tipo}}", data.nvr.tipo) # Añadido: si hay placeholder para tipo de NVR
                replace_text(shape, "{{nvr_direccion}}", data.nvr.direccion)
                replace_text(shape, "{{observaciones_nvr}}", data.nvr.observaciones) # Placeholder común, revisa si aplica aquí
        print("DEBUG: Campos de cantidades/NVR rellenados.", file=sys.stderr)


        # --- Generar slides para cada Cámara ---
        # Buscar el slide que usarás como "prototipo" para cada cámara
        proto = None
        proto_idx = -1 # Para guardar el índice y removerlo después
        print("DEBUG: Buscando slide prototipo con '{{detalle_camara}}' para clonar...", file=sys.stderr)
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                # Verificamos si la forma tiene texto y si contiene el marcador del prototipo
                if shape.has_text_frame:
                    if "{{detalle_camara}}" in shape.text:
                        proto = slide
                        proto_idx = i
                        print(f"DEBUG: Prototipo encontrado en slide con índice: {proto_idx}", file=sys.stderr)
                        break # Encontramos el prototipo, salimos del loop de shapes
            if proto:
                break # Encontramos el prototipo, salimos del loop de slides

        if proto:
            # Remover el slide prototipo de la presentación original
            # Esto es necesario si no quieres que el prototipo aparezca en el resultado final
            prs.slides.remove(proto)
            print(f"DEBUG: Slide prototipo (índice {proto_idx}) removido de la presentación.", file=sys.stderr)

            # Iterar sobre cada cámara en los datos y crear una nueva diapositiva
            print(f"DEBUG: Procesando datos de {len(data.camaras)} cámara(s) para generar slides...", file=sys.stderr)
            for i, cam in enumerate(data.camaras):
                print(f"DEBUG: ---> Procesando Cámara {i+1}/{len(data.camaras)} (Número: {cam.numero})...", file=sys.stderr)

                # Añadir una nueva diapositiva usando el mismo layout del prototipo
                new_sl = prs.slides.add_slide(proto.slide_layout)
                print(f"DEBUG: Slide nuevo creado con el layout del prototipo.", file=sys.stderr)

                # Copiar las formas (shapes) del slide prototipo al nuevo slide
                # Esto copia el contenido visual (texto, cajas, placeholders de imagen, etc.)
                try:
                     print(f"DEBUG: Copiando formas del prototipo al nuevo slide para cámara {cam.numero}...", file=sys.stderr)
                     for s in proto.shapes:
                        # Clonar el elemento XML de la forma y añadirlo al nuevo slide
                        el = s.element
                        new_sl.shapes._spTree.insert_element_before(el.clone(), 'p:extLst')
                     print(f"DEBUG: Formas del prototipo copiadas exitosamente.", file=sys.stderr)
                except Exception as copy_e:
                     print(f"ERROR: Error crítico al copiar formas del prototipo para cámara {cam.numero}: {copy_e}", file=sys.stderr)
                     traceback.print_exc(file=sys.stderr)
                     # Decide aquí si quieres detener la generación o continuar (podría generar slides vacíos)
                     # Por ahora, continuaremos pero registrando el error.


                # Rellenar texto en el nuevo slide de la cámara
                print(f"DEBUG: Rellenando texto en el nuevo slide para cámara {cam.numero}...", file=sys.stderr)
                for shape in new_sl.shapes:
                    # Reemplazar los placeholders específicos de cámara
                    replace_text(shape, "{{detalle_camara}}", f"{cam.numero}. {cam.tipo.upper()}")
                    replace_text(shape, "{{tipo_camara}}", cam.tipo)
                    replace_text(shape, "{{ubicacion_camara}}", cam.ubicacion)
                    replace_text(shape, "{{observaciones_camara}}", cam.observaciones)
                print(f"DEBUG: Texto rellenado para cámara {cam.numero}.", file=sys.stderr)


                # Insertar imagen si hay una URL proporcionada
                print(f"DEBUG: Procesando imagen para cámara {cam.numero}. URL proporcionada: {cam.imagen_url}", file=sys.stderr)
                if cam.imagen_url:
                    img_path = download_image(cam.imagen_url) # Llamamos a la función que descarga (ya tiene logging)
                    if img_path: # Si la descarga fue exitosa (la función devolvió una ruta)
                        try:
                            # === LOGUEAR ANTES DE INSERTAR IMAGEN ===
                            print(f"DEBUG: Intentando insertar imagen desde {img_path} en slide de cámara {cam.numero}...", file=sys.stderr)
                            # Insertar la imagen. Ajusta las coordenadas (left, top) y tamaño (width, height)
                            # para que coincidan con el área donde quieres que aparezca en tu plantilla.
                            # Usamos Inches para medidas.
                            left = Inches(4) # Ejemplo: a 4 pulgadas del borde izquierdo del slide
                            top = Inches(2)  # Ejemplo: a 2 pulgadas del borde superior del slide
                            width = Inches(5) # Ejemplo: la imagen tendrá 5 pulgadas de ancho
                            # height = Inches(3.5) # Opcional: podrías especificar altura si quieres forzar dimensiones

                            new_sl.shapes.add_picture(img_path, left, top, width=width) # Puedes añadir height=height también
                            print(f"DEBUG: Imagen insertada exitosamente en slide de cámara {cam.numero}.", file=sys.stderr)

                            # === LIMPIAR ARCHIVO TEMPORAL DE IMAGEN ===
                            # Es importante eliminar los archivos temporales después de usarlos
                            try:
                                os.remove(img_path)
                                print(f"DEBUG: Archivo temporal de imagen eliminado: {img_path}", file=sys.stderr)
                            except Exception as rm_e:
                                print(f"WARNING: No se pudo eliminar archivo temporal {img_path} después de insertarlo: {rm_e}", file=sys.stderr) # Usar WARNING si la eliminación falla pero el resto funcionó
                        except Exception as img_e:
                            # === LOGUEAR ERRORES AL INSERTAR IMAGEN ===
                            print(f"ERROR: Error al insertar imagen {img_path} en el slide de cámara {cam.numero}: {img_e}", file=sys.stderr)
                            traceback.print_exc(file=sys.stderr)
                    else:
                        # Este caso ya está logueado dentro de download_image, pero un WARNING aquí es útil
                        print(f"WARNING: No se pudo descargar la imagen desde {cam.imagen_url}. No se insertará imagen en slide {cam.numero}.", file=sys.stderr)
                else:
                    print(f"DEBUG: No hay URL de imagen proporcionada para cámara {cam.numero}. Saltando inserción de imagen.", file=sys.stderr)

        else:
            # === LOGUEAR SI NO SE ENCUENTRA EL PROTOTIPO ===
            print("WARNING: No se encontró ningún slide prototipo con '{{detalle_camara}}'. No se generarán slides individuales para cámaras.", file=sys.stderr)
            # Decide si esto es un error crítico. Si es así, podrías lanzar una excepción aquí.
            # raise ValueError("Slide prototipo para cámaras no encontrado en la plantilla.")


        # TODO: Lógica para insertar Mapa si aplica (basado en tu comentario anterior)
        # Si tienes un slide específico para el mapa, necesitarías encontrarlo por un marcador o índice
        # y usar download_image + add_picture similar a como hicimos con las cámaras.
        # Aquí iría esa lógica si la desarrollas.
        # print("DEBUG: Procesando mapa (si aplica)...", file=sys.stderr)
        # ... (tu lógica para el mapa) ...


        # TODO: Lógica para procesar el slide de Cierre si aplica (basado en tu comentario anterior)
        # Si los placeholders de cierre están en un slide específico (ej: el último), procésalo aquí.
        # print("DEBUG: Rellenando campos del slide de cierre (si aplica)...", file=sys.stderr)
        # ... (tu lógica para el cierre) ...


        # === LOGUEAR ANTES DE GUARDAR EL ARCHIVO FINAL ===
        print("DEBUG: Lógica de generación de presentación completada. Procediendo a guardar el archivo final...", file=sys.stderr)

        # Guardar la presentación modificada
        # Creamos un nombre de archivo único en el directorio /tmp/ (escribible en Render)
        out_filename = f"informe_generado_{uuid.uuid4().hex[:8]}.pptx" # Nombre único
        out_path = os.path.join("/tmp", out_filename) # Ruta completa en /tmp/
        try:
            print(f"DEBUG: Intentando guardar presentación final en: {out_path}", file=sys.stderr)
            prs.save(out_path)
            print(f"DEBUG: Presentación guardada exitosamente en: {out_path}", file=sys.stderr)
            return out_path # Devolvemos la ruta del archivo guardado temporalmente
        except Exception as save_e:
            # === LOGUEAR ERRORES AL GUARDAR EL ARCHIVO ===
            print(f"ERROR: Error crítico al guardar el archivo de presentación en {out_path}: {save_e}", file=sys.stderr)
            traceback.print_exc(file=sys.stderr)
            raise save_e # Re-lanzar el error si falla al guardar

    except FileNotFoundError as fnf_e:
         # Capturamos específicamente el error si la plantilla no fue encontrada al inicio
         # Ya se logueó arriba, pero relanzamos para que FastAPI lo capture.
         print(f"ERROR: FileNotFoundError capturado en generar_presentacion: {fnf_e}", file=sys.stderr)
         raise fnf_e # Relanzar la excepción original

    except Exception as e:
        # === LOGUEAR CUALQUIER OTRA EXCEPCIÓN NO MANEJADA EN LA GENERACIÓN ===
        print(f"ERROR: Excepción inesperada durante la generación de la presentación: {e}", file=sys.stderr)
        traceback.print_exc(file=sys.stderr) # Imprimir el rastreo completo
        raise e # Re-lanzar la excepción para que FastAPI la capture


# ---------- FastAPI ---------- #
# Inicializa la aplicación FastAPI. La documentación (/docs, /redoc) está activa por defecto.
app = FastAPI(title="Generador PPT Río de la Plata", description="API para generar presentaciones PowerPoint de forma automatizada a partir de datos JSON.")

# Define el endpoint HTTP POST para generar la presentación
@app.post("/generar-ppt")
async def generar_ppt(data: InformeData):
    # === LOGUEAR EL INICIO DE LA PETICIÓN ===
    print("\n" + "="*50, file=sys.stderr) # Separador en los logs
    print("DEBUG: Recibida petición POST en el endpoint /generar-ppt", file=sys.stderr)
    print("DEBUG: Datos recibidos:", data.model_dump_json(indent=2), file=sys.stderr) # Log de los datos recibidos

    try:
        # Llamar a la función principal para generar la presentación
        ppt_path = generar_presentacion(data)

        # La función generar_presentacion ahora lanza excepciones en caso de errores,
        # por lo que no necesitamos verificar si 'ppt_path is None'.

        # === LOGUEAR ANTES DE DEVOLVER EL ARCHIVO ===
        print(f"DEBUG: Presentación generada exitosamente. Devolviendo archivo desde {ppt_path}", file=sys.stderr)

        # Devolver el archivo PPTX generado como respuesta HTTP
        # filename="informe_generado.pptx" sugiere al navegador cómo nombrar el archivo descargado
        return FileResponse(path=ppt_path, filename="informe_generado.pptx", media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    except FileNotFoundError as fnf_e:
         # Captura específica si la plantilla no se encontró (excepción lanzada en generar_presentacion)
         print(f"ERROR en el endpoint /generar-ppt: Archivo de plantilla no encontrado: {fnf_e}", file=sys.stderr)
         # Devolver una respuesta HTTP 500 con un detalle claro
         raise HTTPException(
             status_code=500,
             detail=f"Error del servidor al generar el PPT: Archivo de plantilla no encontrado en Render. Verifique que '{TEMPLATE_FILENAME}' esté en la raíz del repositorio. Detalle: {fnf_e}"
         )

    except Exception as e:
        # Captura cualquier otro error que haya ocurrido durante la generación
        # El rastreo completo ya fue impreso a sys.stderr en la función generar_presentacion
        print(f"ERROR en el endpoint /generar-ppt: Excepción inesperada: {e}", file=sys.stderr)
        # Devolver una respuesta HTTP 500 indicando un error interno
        raise HTTPException(
            status_code=500,
            detail=f"Error interno del servidor al generar el PPT. Consulte los logs de Render para más detalles sobre la excepción: {e}"
        )
    finally:
        # === LOGUEAR EL FIN DE LA PETICIÓN ===
        print(f"DEBUG: Petición a /generar-ppt finalizada.", file=sys.stderr)
        print("="*50 + "\n", file=sys.stderr) # Separador final
        # Opcional: Limpiar archivos temporales de imagen si no se hizo ya
        # Nota: Los archivos temporales en /tmp suelen limpiarse solos eventualmente,
        # pero si generas muchas imágenes y el proceso vive mucho, podrías querer limpiarlos explícitamente aquí
        # si no se eliminaron inmediatamente después de usarse. La versión actual sí los elimina tras insertarlos.


# Puedes añadir un endpoint raíz simple si quieres, aunque no es estrictamente necesario
# y es normal que dé 404 si no existe.
# @app.get("/")
# def read_root():
#     return {"message": "API Generador PPT funcionando. Vea /docs para la documentación interactiva."}
